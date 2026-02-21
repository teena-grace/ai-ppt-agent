"""
ppt_builder.py — 8 radically different visual templates, each with:
- Unique shapes/geometry
- Unique color application  
- Unique animation choreography
- Unique slide transition

Templates: futuristic, minimalist, artistic, misty, monochrome, hardware, magazine, neon
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, sys, os
sys.path.insert(0, os.path.dirname(__file__))
from anim_engine import AnimSequence, add_transition

W = Inches(13.33); H = Inches(7.5)

def rgb(h):
    h=h.lstrip("#"); return RGBColor(int(h[:2],16),int(h[2:4],16),int(h[4:],16))

def bg(slide, c): f=slide.background.fill; f.solid(); f.fore_color.rgb=c

def rect(slide,l,t,w,h,c):
    s=slide.shapes.add_shape(1,l,t,w,h)
    s.fill.solid(); s.fill.fore_color.rgb=c; s.line.fill.background(); return s

def rrect(slide,l,t,w,h,c,lc=None,lw=1.5):
    """rect with optional border"""
    s=slide.shapes.add_shape(1,l,t,w,h)
    if lc: s.fill.background(); s.line.color.rgb=lc; s.line.width=Pt(lw)
    else: s.fill.solid(); s.fill.fore_color.rgb=c; s.line.fill.background()
    return s

def oval(slide,l,t,w,h,c):
    s=slide.shapes.add_shape(9,l,t,w,h); s.fill.solid()
    s.fill.fore_color.rgb=c; s.line.fill.background(); return s

def tri(slide,l,t,w,h,c):
    s=slide.shapes.add_shape(5,l,t,w,h); s.fill.solid()
    s.fill.fore_color.rgb=c; s.line.fill.background(); return s

def txt(slide,text,l,t,w,h,fn="Calibri",sz=18,c=RGBColor(255,255,255),
        bold=False,italic=False,align=PP_ALIGN.LEFT,wrap=True):
    b=slide.shapes.add_textbox(l,t,w,h); tf=b.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run()
    r.text=text; r.font.name=fn; r.font.size=Pt(sz)
    r.font.color.rgb=c; r.font.bold=bold; r.font.italic=italic; return b

def bpara(tf,text,fn,sz,c,bullet="▸ "):
    p=tf.add_paragraph(); p.alignment=PP_ALIGN.LEFT; p.space_before=Pt(8)
    r=p.add_run(); r.text=bullet+text; r.font.name=fn; r.font.size=Pt(sz); r.font.color.rgb=c

def bullets(slide,pts,l,t,w,h,th,bullet="▸  ",fn=None):
    fn=fn or th["font_body"]
    b=slide.shapes.add_textbox(l,t,w,h); tf=b.text_frame; tf.word_wrap=True
    for i,pt in enumerate(pts):
        if i==0:
            p0=tf.paragraphs[0]; r=p0.add_run(); r.text=bullet+pt
            r.font.name=fn; r.font.size=Pt(13); r.font.color.rgb=rgb(th["body_color"])
        else: bpara(tf,pt,fn,13,rgb(th["body_color"]),bullet)
    return b

def notes(slide,text):
    if text:
        try: slide.notes_slide.notes_text_frame.text=text
        except: pass

# ─────────────────────── TEMPLATE BASE ────────────────────────────────────────
class _Base:
    trans="fade"
    def __init__(self,th): self.th=th

    def _slide(self,prs,bg_col=None):
        s=prs.slides.add_slide(prs.slide_layouts[6])
        if bg_col: bg(s,bg_col)
        else: bg(s,rgb(self.th["bg"]))
        return s, AnimSequence()

    def hero(self,prs,data,idx): return self.content(prs,data,idx)
    def grid(self,prs,data,idx): return self.content(prs,data,idx)
    def timeline(self,prs,data,idx): return self.content(prs,data,idx)
    def numbered(self,prs,data,idx): return self.content(prs,data,idx)

    def _fin(self,slide,seq,data):
        seq.inject(slide); add_transition(slide,self.trans)
        notes(slide,data.get("notes","")); return slide


# ════════════════════════════════════════════════════════
# T1: FUTURISTIC — dark, cyber accent bars, tech circles
# ════════════════════════════════════════════════════════
class Futuristic(_Base):
    trans="push"
    def _chrome(self,slide,seq,idx):
        th=self.th
        bar=rect(slide,Inches(0),Inches(0),Inches(0.42),H,rgb(th["accent"]))
        seq.wipe_in(bar,dur=400)
        tl=rect(slide,Inches(0.42),Inches(0),W-Inches(0.42),Inches(0.05),rgb(th["accent"]))
        seq.wipe_in(tl,dur=500,delay=100)
        c=oval(slide,W-Inches(1.1),Inches(0.08),Inches(0.75),Inches(0.75),rgb(th["accent2"]))
        seq.zoom_in(c,dur=350,delay=200)
        n=txt(slide,f"{idx:02d}",Inches(0.6),Inches(0.12),Inches(1.2),Inches(0.5),
              th["font_body"],11,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=280)

    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._chrome(slide,seq,idx)
        bg_c=oval(slide,Inches(8.5),Inches(-1.8),Inches(6.2),Inches(6.2),rgb(th["header_bg"]))
        seq.zoom_in(bg_c,dur=600,delay=80)
        t=txt(slide,data["title"],Inches(0.6),Inches(1.1),Inches(11),Inches(2.9),
              th["font_head"],54,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=600,delay=280)
        s=txt(slide,data.get("subtitle",""),Inches(0.6),Inches(4.0),Inches(10),Inches(0.8),
              th["font_body"],20,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=400,delay=500)
        d=txt(slide,data.get("detail",""),Inches(0.6),Inches(4.9),Inches(10.5),Inches(2.1),
              th["font_body"],14,rgb(th["muted_color"]),wrap=True)
        seq.fade(d,dur=400,delay=700)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._chrome(slide,seq,idx)
        t=txt(slide,data["title"],Inches(0.6),Inches(0.45),Inches(8.5),Inches(1.2),
              th["font_head"],34,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=240)
        s=txt(slide,data.get("subtitle",""),Inches(0.6),Inches(1.6),Inches(8),Inches(0.5),
              th["font_body"],13,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=350,delay=400)
        bl=bullets(slide,data.get("points",[]),Inches(0.6),Inches(2.25),Inches(6.5),Inches(4.8),th)
        seq.fly_in(bl,dur=500,delay=530)
        card=rect(slide,Inches(7.4),Inches(0.38),Inches(5.5),Inches(6.8),rgb(th["card_bg"]))
        seq.fade(card,dur=400,delay=290)
        cb=rect(slide,Inches(7.4),Inches(0.38),Inches(0.09),Inches(6.8),rgb(th["accent"]))
        seq.wipe_in(cb,dur=400,delay=340)
        lbl=txt(slide,"ANALYSIS",Inches(7.65),Inches(0.72),Inches(5),Inches(0.38),
                th["font_body"],9,rgb(th["accent"]),bold=True)
        seq.appear(lbl,delay=400)
        d=txt(slide,data.get("detail",""),Inches(7.65),Inches(1.2),Inches(5.05),Inches(5.8),
              th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=400,delay=580)
        seq.grow_emphasis(card,scale=102)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._chrome(slide,seq,idx)
        strip=rect(slide,Inches(0.42),Inches(0),W-Inches(0.42),Inches(2.1),rgb(th["header_bg"]))
        seq.fade(strip,dur=400,delay=80)
        t=txt(slide,data["title"],Inches(0.6),Inches(0.35),Inches(12),Inches(1.1),
              th["font_head"],36,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=200)
        pts=data.get("points",[""])
        pos=[(Inches(0.62),Inches(2.28)),(Inches(6.9),Inches(2.28)),(Inches(0.62),Inches(4.98)),(Inches(6.9),Inches(4.98))]
        for i,(pt,(cx,cy)) in enumerate(zip(pts[:4],pos)):
            d=430+i*130
            c=rect(slide,cx,cy,Inches(6.05),Inches(2.38),rgb(th["card_bg"]))
            seq.zoom_in(c,dur=400,delay=d)
            bdg=oval(slide,cx+Inches(0.14),cy+Inches(0.13),Inches(0.54),Inches(0.54),rgb(th["accent"]))
            seq.zoom_in(bdg,dur=250,delay=d+80); seq.grow_emphasis(bdg,scale=115)
            bn=txt(slide,f"0{i+1}",cx+Inches(0.14),cy+Inches(0.13),Inches(0.54),Inches(0.54),
                   th["font_body"],11,rgb(th["bg"]),bold=True,align=PP_ALIGN.CENTER)
            seq.appear(bn,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.84),cy+Inches(0.08),Inches(5.05),Inches(2.18),
                     th["font_body"],12,rgb(th["body_color"]),wrap=True)
            seq.fade(ptxt,dur=350,delay=d+150)
        return self._fin(slide,seq,data)

    def timeline(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._chrome(slide,seq,idx)
        t=txt(slide,data["title"],Inches(0.6),Inches(0.38),Inches(12),Inches(1.1),
              th["font_head"],36,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=200)
        tl=rect(slide,Inches(0.6),Inches(3.48),Inches(12.45),Inches(0.08),rgb(th["accent"]))
        seq.wipe_in(tl,dur=600,delay=420)
        pts=data.get("points",[]); n=min(len(pts),4); sp=Inches(12.4)/max(n,1)
        for i,pt in enumerate(pts[:4]):
            cx=Inches(0.6)+i*sp+sp/2-Inches(1.5); d=520+i*160
            dot=oval(slide,cx+Inches(1.14),Inches(3.2),Inches(0.6),Inches(0.6),rgb(th["accent"]))
            seq.zoom_in(dot,dur=300,delay=d); seq.grow_emphasis(dot,scale=120)
            sn=txt(slide,str(i+1),cx+Inches(1.14),Inches(3.2),Inches(0.6),Inches(0.6),
                   th["font_head"],13,rgb(th["bg"]),bold=True,align=PP_ALIGN.CENTER)
            seq.appear(sn,delay=d+40)
            card=rect(slide,cx,Inches(3.98),Inches(3.04),Inches(2.98),rgb(th["card_bg"]))
            seq.fade(card,dur=350,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.12),Inches(4.12),Inches(2.78),Inches(2.72),
                     th["font_body"],12,rgb(th["body_color"]),wrap=True)
            seq.fly_in(ptxt,dur=350,delay=d+160)
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._chrome(slide,seq,idx)
        panel=rect(slide,Inches(0.42),Inches(0),Inches(5.25),H,rgb(th["header_bg"]))
        seq.fade(panel,dur=400,delay=80)
        sep=rect(slide,Inches(5.67),Inches(0),Inches(0.08),H,rgb(th["accent"]))
        seq.wipe_in(sep,dur=450,delay=120)
        t=txt(slide,data["title"],Inches(0.62),Inches(0.68),Inches(4.8),Inches(1.5),
              th["font_head"],30,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=240)
        s=txt(slide,data.get("subtitle",""),Inches(0.62),Inches(2.22),Inches(4.8),Inches(0.55),
              th["font_body"],13,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=350,delay=400)
        d=txt(slide,data.get("detail",""),Inches(0.62),Inches(2.95),Inches(4.8),Inches(4.1),
              th["font_body"],12,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=400,delay=530)
        for i,pt in enumerate(data.get("points",[])[:4]):
            y=Inches(0.52+i*1.65); dd=290+i*160
            c=oval(slide,Inches(6.08),y,Inches(0.68),Inches(0.68),rgb(th["accent"]))
            seq.zoom_in(c,dur=300,delay=dd)
            nt=txt(slide,str(i+1),Inches(6.08),y,Inches(0.68),Inches(0.68),
                   th["font_head"],16,rgb(th["bg"]),bold=True,align=PP_ALIGN.CENTER)
            seq.appear(nt,delay=dd+30)
            ptxt=txt(slide,pt,Inches(6.9),y,Inches(6.1),Inches(1.5),
                     th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.sweep_from_left(ptxt,delay=dd+80,dur=600)
        return self._fin(slide,seq,data)


# ════════════════════════════════════════════════════════
# T2: MINIMALIST — white space, hairlines, type-forward
# ════════════════════════════════════════════════════════
class Minimalist(_Base):
    trans="fade"
    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        hr=rect(slide,Inches(1.0),Inches(1.45),Inches(3.5),Inches(0.025),rgb(th["accent"]))
        seq.wipe_in(hr,dur=600)
        wm=txt(slide,f"{idx:02d}",Inches(9.5),Inches(4.5),Inches(3.5),Inches(3.0),
               th["font_head"],120,rgb(th["muted_color"]),bold=True,align=PP_ALIGN.RIGHT)
        seq.fade(wm,dur=800,delay=100)
        t=txt(slide,data["title"],Inches(1.0),Inches(1.75),Inches(9.5),Inches(2.8),
              th["font_head"],54,rgb(th["title_color"]),bold=True)
        seq.fade(t,dur=600,delay=200)
        s=txt(slide,data.get("subtitle",""),Inches(1.0),Inches(4.5),Inches(9),Inches(0.7),
              th["font_body"],16,rgb(th["muted_color"]),italic=True)
        seq.fade(s,dur=500,delay=450)
        d=txt(slide,data.get("detail",""),Inches(1.0),Inches(5.3),Inches(9.5),Inches(1.8),
              th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=500,delay=650)
        br=rect(slide,Inches(1.0),Inches(6.98),W-Inches(2.0),Inches(0.025),rgb(th["muted_color"]))
        seq.wipe_in(br,dur=600,delay=800)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        vl=rect(slide,Inches(1.0),Inches(0.8),Inches(0.025),Inches(6.0),rgb(th["accent"]))
        seq.wipe_in(vl,dur=500)
        n=txt(slide,f"{idx:02d}",Inches(0.38),Inches(0.38),Inches(0.5),Inches(0.5),
              th["font_body"],10,rgb(th["muted_color"]))
        seq.appear(n,delay=100)
        t=txt(slide,data["title"],Inches(1.2),Inches(0.65),Inches(9.5),Inches(1.3),
              th["font_head"],36,rgb(th["title_color"]),bold=True)
        seq.fade(t,dur=500,delay=200)
        s=txt(slide,data.get("subtitle",""),Inches(1.2),Inches(1.9),Inches(9.5),Inches(0.5),
              th["font_body"],13,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=400,delay=370)
        bl=bullets(slide,data.get("points",[]),Inches(1.2),Inches(2.55),Inches(6.5),Inches(4.5),th,bullet="— ")
        seq.fade(bl,dur=450,delay=490)
        db_t=rect(slide,Inches(8.5),Inches(0.65),Inches(4.35),Inches(0.04),rgb(th["accent"]))
        seq.wipe_in(db_t,dur=400,delay=310)
        d=txt(slide,data.get("detail",""),Inches(8.5),Inches(0.88),Inches(4.35),Inches(6.3),
              th["font_body"],14,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=500,delay=510)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        hr=rect(slide,Inches(1.0),Inches(0.98),Inches(11.5),Inches(0.025),rgb(th["accent"]))
        seq.wipe_in(hr,dur=600)
        t=txt(slide,data["title"],Inches(1.0),Inches(0.22),Inches(10),Inches(0.8),
              th["font_head"],32,rgb(th["title_color"]),bold=True)
        seq.fade(t,dur=500,delay=150)
        pts=data.get("points",[""])
        pos=[(Inches(1.0),Inches(1.28)),(Inches(7.3),Inches(1.28)),(Inches(1.0),Inches(4.18)),(Inches(7.3),Inches(4.18))]
        for i,(pt,(cx,cy)) in enumerate(zip(pts[:4],pos)):
            d=340+i*120
            cb=rect(slide,cx,cy,Inches(0.03),Inches(2.8),rgb(th["accent"]))
            seq.wipe_in(cb,dur=350,delay=d)
            num=txt(slide,f"0{i+1}",cx+Inches(0.14),cy,Inches(0.8),Inches(0.5),
                    th["font_head"],18,rgb(th["accent"]),bold=True)
            seq.fade(num,dur=250,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.14),cy+Inches(0.44),Inches(6.0),Inches(2.3),
                     th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.fade(ptxt,dur=350,delay=d+150)
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        vl=rect(slide,Inches(1.0),Inches(0.48),Inches(0.025),Inches(6.7),rgb(th["muted_color"]))
        seq.wipe_in(vl,dur=600)
        t=txt(slide,data["title"],Inches(1.3),Inches(0.38),Inches(11.5),Inches(1.2),
              th["font_head"],38,rgb(th["title_color"]),bold=True)
        seq.fade(t,dur=500,delay=150)
        ac=rect(slide,Inches(1.3),Inches(1.48),Inches(2.5),Inches(0.04),rgb(th["accent"]))
        seq.wipe_in(ac,dur=400,delay=280)
        for i,pt in enumerate(data.get("points",[])[:4]):
            y=Inches(1.78+i*1.35); d=340+i*150
            num=txt(slide,f"0{i+1}",Inches(1.3),y,Inches(0.7),Inches(0.7),
                    th["font_head"],22,rgb(th["accent"]),bold=True)
            seq.fade(num,dur=300,delay=d)
            hl=rect(slide,Inches(2.1),y+Inches(0.32),Inches(0.5),Inches(0.02),rgb(th["muted_color"]))
            seq.wipe_in(hl,dur=300,delay=d+80)
            ptxt=txt(slide,pt,Inches(2.78),y,Inches(10.1),Inches(1.2),
                     th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.fade(ptxt,dur=400,delay=d+120)
        return self._fin(slide,seq,data)

    def timeline(self,prs,data,idx): return self.numbered(prs,data,idx)


# ════════════════════════════════════════════════════════
# T3: ARTISTIC — overlapping geometry, editorial, expressive
# ════════════════════════════════════════════════════════
class Artistic(_Base):
    trans="zoom"
    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        b1=oval(slide,Inches(-1),Inches(-1),Inches(6),Inches(6),rgb(th["header_bg"]))
        seq.zoom_in(b1,dur=700)
        b2=oval(slide,Inches(9),Inches(3.5),Inches(5),Inches(5),rgb(th["accent2"]))
        seq.zoom_in(b2,dur=700,delay=100)
        b3=rect(slide,Inches(5),Inches(-0.5),Inches(8.33),Inches(0.5),rgb(th["accent"]))
        seq.wipe_in(b3,dur=500,delay=200)
        t3=tri(slide,Inches(11),Inches(5),Inches(2),Inches(2.5),rgb(th["accent"]))
        seq.zoom_in(t3,dur=400,delay=300); seq.spin_emphasis(t3,degrees=45)
        n=txt(slide,f"{idx:02d}",Inches(0.5),Inches(0.28),Inches(1.5),Inches(0.7),
              th["font_head"],14,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=350)
        t=txt(slide,data["title"],Inches(0.6),Inches(1.4),Inches(9.5),Inches(3.0),
              th["font_head"],58,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=600,delay=280,from_dir="left")
        s=txt(slide,data.get("subtitle",""),Inches(0.6),Inches(4.35),Inches(9),Inches(0.8),
              th["font_body"],18,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=400,delay=550)
        d=txt(slide,data.get("detail",""),Inches(0.6),Inches(5.25),Inches(10),Inches(1.8),
              th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=400,delay=720)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        rp=rect(slide,Inches(7.8),Inches(0),Inches(5.53),H,rgb(th["header_bg"]))
        seq.fade(rp,dur=500)
        dt=tri(slide,Inches(7.3),Inches(0),Inches(1.5),H,rgb(th["accent2"]))
        seq.zoom_in(dt,dur=500,delay=100)
        ac=oval(slide,Inches(9.5),Inches(5.2),Inches(3.5),Inches(3.5),rgb(th["accent"]))
        seq.zoom_in(ac,dur=400,delay=200)
        n=txt(slide,f"{idx:02d}",Inches(0.3),Inches(0.2),Inches(1),Inches(0.6),
              th["font_head"],13,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=250)
        t=txt(slide,data["title"],Inches(0.3),Inches(0.65),Inches(7.2),Inches(1.4),
              th["font_head"],34,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=300,from_dir="left")
        s=txt(slide,data.get("subtitle",""),Inches(0.3),Inches(1.98),Inches(7.0),Inches(0.5),
              th["font_body"],14,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=350,delay=450)
        bl=bullets(slide,data.get("points",[]),Inches(0.3),Inches(2.65),Inches(6.8),Inches(4.5),th,bullet="● ")
        seq.sweep_from_left(bl,delay=550,dur=600)
        d=txt(slide,data.get("detail",""),Inches(8.3),Inches(1.5),Inches(4.7),Inches(5.6),
              th["font_body"],13,rgb(th["title_color"]),wrap=True)
        seq.fade(d,dur=400,delay=480)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        dc=oval(slide,Inches(10.5),Inches(-1),Inches(4),Inches(4),rgb(th["header_bg"]))
        seq.zoom_in(dc,dur=600)
        t=txt(slide,data["title"],Inches(0.5),Inches(0.28),Inches(9.5),Inches(1.2),
              th["font_head"],38,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=200,from_dir="left")
        pts=data.get("points",[""])
        for i,pt in enumerate(pts[:4]):
            row=i//2; col=i%2
            cx=Inches(0.5)+col*Inches(6.5); cy=Inches(1.78)+row*Inches(2.8)
            d=400+i*140
            dm=tri(slide,cx,cy,Inches(0.78),Inches(0.78),rgb(th["accent"]))
            seq.zoom_in(dm,dur=350,delay=d); seq.spin_emphasis(dm,degrees=45)
            c=rect(slide,cx+Inches(0.48),cy,Inches(5.72),Inches(2.55),rgb(th["card_bg"]))
            seq.fade(c,dur=350,delay=d+50)
            num=txt(slide,f"0{i+1}",cx+Inches(0.58),cy+Inches(0.08),Inches(0.6),Inches(0.45),
                    th["font_head"],16,rgb(th["accent"]),bold=True)
            seq.appear(num,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.58),cy+Inches(0.5),Inches(5.3),Inches(2.0),
                     th["font_body"],12,rgb(th["body_color"]),wrap=True)
            seq.fade(ptxt,dur=350,delay=d+150)
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        diag=rect(slide,Inches(0),Inches(5.5),W,Inches(2.5),rgb(th["header_bg"]))
        seq.fade(diag,dur=500)
        bc=oval(slide,Inches(-1),Inches(-1.5),Inches(5),Inches(5),rgb(th["accent2"]))
        seq.zoom_in(bc,dur=600)
        t=txt(slide,data["title"],Inches(0.5),Inches(0.48),Inches(12),Inches(1.2),
              th["font_head"],36,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=250,from_dir="left")
        for i,pt in enumerate(data.get("points",[])[:4]):
            y=Inches(1.88+i*1.2); d=370+i*150
            c=oval(slide,Inches(0.4),y,Inches(0.6),Inches(0.6),rgb(th["accent"]))
            seq.zoom_in(c,dur=300,delay=d); seq.pulse(c)
            nt=txt(slide,str(i+1),Inches(0.4),y,Inches(0.6),Inches(0.6),
                   th["font_head"],14,rgb(th["bg"]),bold=True,align=PP_ALIGN.CENTER)
            seq.appear(nt,delay=d+40)
            ptxt=txt(slide,pt,Inches(1.2),y,Inches(11.5),Inches(1.1),
                     th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.sweep_from_left(ptxt,delay=d+80,dur=600)
        return self._fin(slide,seq,data)

    def timeline(self,prs,data,idx): return self.numbered(prs,data,idx)


# ════════════════════════════════════════════════════════
# T4: MISTY PEAKS — mountain silhouettes, atmospheric
# ════════════════════════════════════════════════════════
class MistyPeaks(_Base):
    trans="dissolve"
    def _mountains(self,slide,seq,th):
        m1=tri(slide,Inches(-0.5),Inches(1.5),Inches(5),Inches(6.5),rgb(th["header_bg"]))
        seq.fade(m1,dur=800)
        m2=tri(slide,Inches(3),Inches(2.5),Inches(5.5),Inches(5.5),rgb(th["header_bg"]))
        seq.fade(m2,dur=800,delay=80)
        m3=tri(slide,Inches(6),Inches(3.0),Inches(5),Inches(5.0),rgb(th["card_bg"]))
        seq.fade(m3,dur=700,delay=200)
        m4=tri(slide,Inches(9),Inches(2.0),Inches(5.5),Inches(6.0),rgb(th["card_bg"]))
        seq.fade(m4,dur=700,delay=280)
        base=rect(slide,Inches(0),Inches(5.5),W,Inches(2.0),rgb(th["accent2"]))
        seq.fade(base,dur=600,delay=350)

    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._mountains(slide,seq,th)
        mist=rect(slide,Inches(0),Inches(3.8),W,Inches(1.5),rgb(th["bg"]))
        seq.fade(mist,dur=600,delay=450)
        n=txt(slide,f"{idx:02d}",Inches(0.5),Inches(0.28),Inches(1.2),Inches(0.5),
              th["font_body"],11,rgb(th["muted_color"]))
        seq.appear(n,delay=500)
        t=txt(slide,data["title"],Inches(0.5),Inches(0.48),Inches(12),Inches(2.5),
              th["font_head"],54,rgb(th["title_color"]),bold=True)
        seq.float_up(t,dur=700,delay=500)
        s=txt(slide,data.get("subtitle",""),Inches(0.5),Inches(2.98),Inches(11),Inches(0.8),
              th["font_body"],18,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=500,delay=750)
        d=txt(slide,data.get("detail",""),Inches(0.5),Inches(3.88),Inches(11),Inches(1.5),
              th["font_body"],13,rgb(th["muted_color"]),wrap=True)
        seq.fade(d,dur=500,delay=950)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        m1=tri(slide,Inches(7.5),Inches(1.5),Inches(4.5),Inches(6.0),rgb(th["header_bg"]))
        seq.fade(m1,dur=700)
        m2=tri(slide,Inches(10),Inches(2.5),Inches(4.0),Inches(5.0),rgb(th["card_bg"]))
        seq.fade(m2,dur=700,delay=100)
        base=rect(slide,Inches(0),Inches(6.2),W,Inches(1.3),rgb(th["accent2"]))
        seq.fade(base,dur=600,delay=200)
        n=txt(slide,f"{idx:02d}",Inches(0.4),Inches(0.28),Inches(1),Inches(0.45),
              th["font_body"],10,rgb(th["muted_color"]))
        seq.appear(n,delay=300)
        t=txt(slide,data["title"],Inches(0.4),Inches(0.58),Inches(7.5),Inches(1.4),
              th["font_head"],34,rgb(th["title_color"]),bold=True)
        seq.float_up(t,dur=600,delay=350)
        s=txt(slide,data.get("subtitle",""),Inches(0.4),Inches(1.92),Inches(7.5),Inches(0.5),
              th["font_body"],13,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=400,delay=520)
        hl=rect(slide,Inches(0.4),Inches(2.52),Inches(4.0),Inches(0.025),rgb(th["accent"]))
        seq.wipe_in(hl,dur=400,delay=550)
        bl=bullets(slide,data.get("points",[]),Inches(0.4),Inches(2.72),Inches(6.8),Inches(3.5),th,bullet="· ")
        seq.fade(bl,dur=500,delay=650)
        d=txt(slide,data.get("detail",""),Inches(0.4),Inches(5.88),Inches(7.0),Inches(1.2),
              th["font_body"],12,rgb(th["muted_color"]),wrap=True)
        seq.fade(d,dur=400,delay=850)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        base=rect(slide,Inches(0),Inches(6.48),W,Inches(1.02),rgb(th["accent2"]))
        seq.fade(base,dur=500)
        m1=tri(slide,Inches(11),Inches(2),Inches(3),Inches(5.5),rgb(th["header_bg"]))
        seq.fade(m1,dur=600,delay=50)
        t=txt(slide,data["title"],Inches(0.4),Inches(0.22),Inches(10),Inches(0.9),
              th["font_head"],32,rgb(th["title_color"]),bold=True)
        seq.float_up(t,dur=600,delay=200)
        hl=rect(slide,Inches(0.4),Inches(1.12),Inches(3.5),Inches(0.03),rgb(th["accent"]))
        seq.wipe_in(hl,dur=400,delay=380)
        pts=data.get("points",[""])
        pos=[(Inches(0.4),Inches(1.38)),(Inches(6.7),Inches(1.38)),(Inches(0.4),Inches(3.98)),(Inches(6.7),Inches(3.98))]
        for i,(pt,(cx,cy)) in enumerate(zip(pts[:4],pos)):
            d=440+i*130
            c=rect(slide,cx,cy,Inches(6.0),Inches(2.45),rgb(th["card_bg"]))
            seq.fade(c,dur=400,delay=d)
            tb=rect(slide,cx,cy,Inches(6.0),Inches(0.05),rgb(th["accent"]))
            seq.wipe_in(tb,dur=350,delay=d+60)
            num=txt(slide,f"0{i+1}",cx+Inches(0.14),cy+Inches(0.1),Inches(0.7),Inches(0.5),
                    th["font_head"],18,rgb(th["accent"]),bold=True)
            seq.appear(num,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.14),cy+Inches(0.52),Inches(5.65),Inches(1.82),
                     th["font_body"],12,rgb(th["body_color"]),wrap=True)
            seq.float_up(ptxt,dur=400,delay=d+150)
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx): return self.grid(prs,data,idx)
    def timeline(self,prs,data,idx): return self.content(prs,data,idx)


# ════════════════════════════════════════════════════════
# T5: MONOCHROME — single hue, stark, heavy type
# ════════════════════════════════════════════════════════
class Monochrome(_Base):
    trans="wipe"
    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        half=rect(slide,Inches(0),Inches(0),Inches(6.5),H,rgb(th["header_bg"]))
        seq.wipe_in(half,dur=600)
        big_n=txt(slide,f"{idx:02d}",Inches(0.3),Inches(0.2),Inches(2.5),Inches(1.2),
                  th["font_head"],72,rgb(th["accent"]),bold=True)
        seq.wipe_in(big_n,dur=400,delay=200)
        t=txt(slide,data["title"],Inches(0.3),Inches(1.28),Inches(6.0),Inches(3.5),
              th["font_head"],52,rgb(th["title_color"]),bold=True)
        seq.wipe_in(t,dur=500,delay=300)
        s=txt(slide,data.get("subtitle",""),Inches(0.3),Inches(4.78),Inches(6.0),Inches(0.75),
              th["font_body"],15,rgb(th["accent"]),italic=True)
        seq.wipe_in(s,dur=400,delay=500)
        vd=rect(slide,Inches(6.5),Inches(0),Inches(0.06),H,rgb(th["accent"]))
        seq.wipe_in(vd,dur=400,delay=350)
        d=txt(slide,data.get("detail",""),Inches(6.8),Inches(1.5),Inches(6.2),Inches(5.5),
              th["font_body"],16,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=600,delay=600)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        tb=rect(slide,Inches(0),Inches(0),W,Inches(1.5),rgb(th["header_bg"]))
        seq.wipe_in(tb,dur=600)
        n=txt(slide,f"{idx:02d}",Inches(0.3),Inches(0.08),Inches(1.2),Inches(0.7),
              th["font_head"],28,rgb(th["accent"]),bold=True)
        seq.wipe_in(n,dur=300,delay=150)
        t=txt(slide,data["title"],Inches(1.5),Inches(0.08),Inches(11.3),Inches(1.2),
              th["font_head"],34,rgb(th["title_color"]),bold=True)
        seq.wipe_in(t,dur=500,delay=250)
        for i,pt in enumerate(data.get("points",[])[:4]):
            y=Inches(1.68+i*1.3); d=390+i*150
            nb=rect(slide,Inches(0.3),y,Inches(0.8),Inches(0.9),rgb(th["accent"]))
            seq.wipe_in(nb,dur=300,delay=d)
            num=txt(slide,f"{i+1}",Inches(0.3),y,Inches(0.8),Inches(0.9),
                    th["font_head"],22,rgb(th["bg"]),bold=True,align=PP_ALIGN.CENTER)
            seq.appear(num,delay=d+30)
            ptxt=txt(slide,pt,Inches(1.28),y+Inches(0.05),Inches(7.0),Inches(1.15),
                     th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.wipe_in(ptxt,dur=400,delay=d+80)
        vd2=rect(slide,Inches(8.5),Inches(1.5),Inches(0.06),Inches(5.7),rgb(th["accent"]))
        seq.wipe_in(vd2,dur=450,delay=350)
        d_txt=txt(slide,data.get("detail",""),Inches(8.7),Inches(1.58),Inches(4.3),Inches(5.6),
                  th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.fade(d_txt,dur=500,delay=600)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        b1=rect(slide,Inches(0),Inches(0),Inches(6.5),Inches(3.58),rgb(th["header_bg"]))
        seq.wipe_in(b1,dur=500)
        b2=rect(slide,Inches(6.5),Inches(3.58),Inches(6.83),Inches(3.92),rgb(th["header_bg"]))
        seq.wipe_in(b2,dur=500,delay=100)
        t=txt(slide,data["title"],Inches(0.3),Inches(0.08),Inches(11),Inches(0.9),
              th["font_head"],34,rgb(th["title_color"]),bold=True)
        seq.wipe_in(t,dur=500,delay=200)
        pts=data.get("points",[""])
        quadrant_tc=[rgb(th["title_color"]),rgb(th["body_color"]),rgb(th["body_color"]),rgb(th["title_color"])]
        for i,pt in enumerate(pts[:4]):
            cx=Inches(0.3) if i%2==0 else Inches(6.8)
            cy=Inches(1.0) if i<2 else Inches(4.48)
            d=370+i*120
            num=txt(slide,f"0{i+1}",cx,cy-Inches(0.38),Inches(1.5),Inches(0.5),
                    th["font_head"],18,rgb(th["accent"]),bold=True)
            seq.wipe_in(num,dur=300,delay=d)
            ptxt=txt(slide,pt,cx,cy,Inches(6.0),Inches(2.3),
                     th["font_body"],13,quadrant_tc[i],wrap=True)
            seq.wipe_in(ptxt,dur=400,delay=d+100)
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx): return self.content(prs,data,idx)
    def timeline(self,prs,data,idx): return self.content(prs,data,idx)


# ════════════════════════════════════════════════════════
# T6: HARDWARE — dark blueprint, circuit grid, code font
# ════════════════════════════════════════════════════════
class Hardware(_Base):
    trans="cut"
    def _grid(self,slide,seq):
        th=self.th
        for i in range(5):
            l=rect(slide,Inches(1+i*2.3),Inches(0),Inches(0.02),H,rgb(th["muted_color"]))
            seq.appear(l,delay=50+i*25)
        for j in range(3):
            l=rect(slide,Inches(0),Inches(1.2+j*2.1),W,Inches(0.02),rgb(th["muted_color"]))
            seq.appear(l,delay=75+j*25)
        for (lx,ly) in [(Inches(0.15),Inches(0.15)),(W-Inches(0.65),Inches(0.15)),
                        (Inches(0.15),H-Inches(0.5)),(W-Inches(0.65),H-Inches(0.5))]:
            b=rect(slide,lx,ly,Inches(0.48),Inches(0.04),rgb(th["accent"]))
            seq.appear(b,delay=40)
            b2=rect(slide,lx,ly,Inches(0.04),Inches(0.4),rgb(th["accent"]))
            seq.appear(b2,delay=40)

    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._grid(slide,seq)
        cpu=rect(slide,Inches(0.38),Inches(1.5),Inches(0.15),Inches(4.5),rgb(th["accent"]))
        seq.appear(cpu,delay=100)
        for y_off in [0,2.2,4.4]:
            t_=rect(slide,Inches(0.53),Inches(1.5+y_off),Inches(0.6),Inches(0.1),rgb(th["accent"]))
            seq.appear(t_,delay=120)
        st=rect(slide,Inches(1.2),Inches(0.7),Inches(9),Inches(0.04),rgb(th["accent"]))
        seq.wipe_in(st,dur=600,delay=200)
        n=txt(slide,f"/{idx:02d}",Inches(1.2),Inches(0.12),Inches(2),Inches(0.6),
              th["font_body"],12,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=200)
        t=txt(slide,data["title"],Inches(1.2),Inches(0.98),Inches(11.8),Inches(2.8),
              th["font_head"],52,rgb(th["title_color"]),bold=True)
        seq.appear(t,delay=350)
        s=txt(slide,data.get("subtitle",""),Inches(1.2),Inches(3.78),Inches(10),Inches(0.7),
              th["font_body"],16,rgb(th["accent"]))
        seq.appear(s,delay=480)
        d=txt(slide,data.get("detail",""),Inches(1.2),Inches(4.68),Inches(11.8),Inches(2.3),
              th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.appear(d,delay=620)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._grid(slide,seq)
        hdr=rect(slide,Inches(0),Inches(0),W,Inches(1.4),rgb(th["header_bg"]))
        seq.appear(hdr,delay=50)
        n=txt(slide,f"[{idx:02d}]",Inches(0.3),Inches(0.1),Inches(1.5),Inches(0.7),
              th["font_body"],16,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=100)
        t=txt(slide,data["title"],Inches(1.8),Inches(0.1),Inches(10),Inches(1.1),
              th["font_head"],30,rgb(th["title_color"]),bold=True)
        seq.appear(t,delay=150)
        lb=rect(slide,Inches(0),Inches(1.4),Inches(5.5),Inches(6.1),rgb(th["card_bg"]))
        seq.appear(lb,delay=200)
        lbl=txt(slide,"// SPECS",Inches(0.2),Inches(1.48),Inches(5),Inches(0.45),
                th["font_body"],10,rgb(th["accent"]),bold=True)
        seq.appear(lbl,delay=250)
        sp=rect(slide,Inches(0.2),Inches(1.93),Inches(5.0),Inches(0.04),rgb(th["accent"]))
        seq.wipe_in(sp,dur=400,delay=280)
        bl=bullets(slide,data.get("points",[]),Inches(0.2),Inches(2.08),Inches(5.0),Inches(5.0),th,bullet="> ",fn="Courier New")
        seq.appear(bl,delay=350)
        lbl2=txt(slide,"// ANALYSIS",Inches(5.7),Inches(1.48),Inches(7.3),Inches(0.45),
                 th["font_body"],10,rgb(th["accent"]),bold=True)
        seq.appear(lbl2,delay=280)
        sp2=rect(slide,Inches(5.7),Inches(1.93),Inches(7.3),Inches(0.04),rgb(th["accent"]))
        seq.wipe_in(sp2,dur=400,delay=300)
        d=txt(slide,data.get("detail",""),Inches(5.7),Inches(2.08),Inches(7.3),Inches(5.0),
              th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.appear(d,delay=400)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        self._grid(slide,seq)
        hdr=rect(slide,Inches(0),Inches(0),W,Inches(1.5),rgb(th["header_bg"]))
        seq.appear(hdr,delay=50)
        n=txt(slide,f"[{idx:02d}]",Inches(0.3),Inches(0.1),Inches(1.5),Inches(0.7),
              th["font_body"],16,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=100)
        t=txt(slide,data["title"],Inches(1.8),Inches(0.12),Inches(10.5),Inches(1.05),
              th["font_head"],30,rgb(th["title_color"]),bold=True)
        seq.appear(t,delay=150)
        pts=data.get("points",[""])
        pos=[(Inches(0.3),Inches(1.68)),(Inches(6.7),Inches(1.68)),(Inches(0.3),Inches(4.58)),(Inches(6.7),Inches(4.58))]
        for i,(pt,(cx,cy)) in enumerate(zip(pts[:4],pos)):
            d=290+i*120
            box=rect(slide,cx,cy,Inches(6.2),Inches(2.65),rgb(th["card_bg"]))
            seq.appear(box,delay=d)
            tb=rect(slide,cx,cy,Inches(6.2),Inches(0.05),rgb(th["accent"]))
            seq.wipe_in(tb,dur=350,delay=d+50)
            lbl=txt(slide,f"MODULE_{i+1:02d}",cx+Inches(0.14),cy+Inches(0.08),Inches(3),Inches(0.4),
                    th["font_body"],9,rgb(th["accent"]),bold=True)
            seq.appear(lbl,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.14),cy+Inches(0.48),"Courier New" if False else th["font_body"],
                     0 if False else 11,rgb(th["body_color"]),wrap=True) if False else \
                 txt(slide,pt,cx+Inches(0.14),cy+Inches(0.48),Inches(5.85),Inches(2.0),
                     "Courier New",11,rgb(th["body_color"]),wrap=True)
            seq.appear(ptxt,delay=d+120)
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx): return self.grid(prs,data,idx)
    def timeline(self,prs,data,idx): return self.content(prs,data,idx)


# ════════════════════════════════════════════════════════
# T7: MAGAZINE — editorial, bold, full-bleed asymmetry
# ════════════════════════════════════════════════════════
class Magazine(_Base):
    trans="cover"
    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        rp=rect(slide,Inches(7.5),Inches(0),Inches(5.83),H,rgb(th["header_bg"]))
        seq.fade(rp,dur=600)
        bn=txt(slide,f"{idx:02d}",Inches(7.6),Inches(-0.5),Inches(5.5),Inches(5),
               th["font_head"],200,rgb(th["bg"]),bold=True)
        seq.fade(bn,dur=500,delay=100)
        tag=txt(slide,"ISSUE // "+str(idx).zfill(2),Inches(0.4),Inches(0.32),Inches(7),Inches(0.5),
                th["font_body"],11,rgb(th["accent"]),bold=True)
        seq.appear(tag,delay=150)
        t=txt(slide,data["title"],Inches(0.4),Inches(0.98),Inches(6.8),Inches(3.8),
              th["font_head"],60,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=600,delay=200,from_dir="bottom")
        rule=rect(slide,Inches(0.4),Inches(4.82),Inches(3.5),Inches(0.08),rgb(th["accent"]))
        seq.wipe_in(rule,dur=400,delay=480)
        s=txt(slide,data.get("subtitle",""),Inches(0.4),Inches(4.98),Inches(6.8),Inches(0.7),
              th["font_body"],16,rgb(th["muted_color"]),italic=True)
        seq.fade(s,dur=400,delay=620)
        d=txt(slide,data.get("detail",""),Inches(7.7),Inches(5.48),Inches(5.3),Inches(1.8),
              th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=400,delay=700)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        bar=rect(slide,Inches(0),Inches(0),Inches(0.6),H,rgb(th["accent"]))
        seq.wipe_in(bar,dur=400)
        t=txt(slide,data["title"],Inches(0.8),Inches(0.12),Inches(11),Inches(1.5),
              th["font_head"],42,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=200,from_dir="left")
        bot=rect(slide,Inches(0.6),Inches(1.62),W-Inches(0.6),Inches(0.07),rgb(th["accent"]))
        seq.wipe_in(bot,dur=500,delay=380)
        s=txt(slide,data.get("subtitle",""),Inches(0.8),Inches(1.75),Inches(11),Inches(0.5),
              th["font_body"],14,rgb(th["accent"]),italic=True)
        seq.fade(s,delay=480)
        pts=data.get("points",[]); lp=pts[:2]; rp=pts[2:4]
        for i,pt in enumerate(lp):
            b=txt(slide,f"↳  {pt}",Inches(0.8),Inches(2.48+i*1.7),Inches(6.0),Inches(1.6),
                  th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.sweep_from_left(b,delay=550+i*150,dur=600)
        for i,pt in enumerate(rp):
            b=txt(slide,f"↳  {pt}",Inches(7.0),Inches(2.48+i*1.7),Inches(6.0),Inches(1.6),
                  th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.sweep_from_right(b,delay=550+i*150,dur=600)
        ib=rect(slide,Inches(0.8),Inches(5.82),Inches(12.2),Inches(1.42),rgb(th["card_bg"]))
        seq.zoom_in(ib,dur=400,delay=800); seq.grow_emphasis(ib,scale=102)
        d=txt(slide,data.get("detail",""),Inches(1.0),Inches(5.95),Inches(12.0),Inches(1.1),
              th["font_body"],13,rgb(th["body_color"]),wrap=True)
        seq.fade(d,delay=900)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        wm=txt(slide,str(idx),Inches(8.5),Inches(-0.8),Inches(5.0),Inches(6.0),
               th["font_head"],220,rgb(th["header_bg"]),bold=True)
        seq.fade(wm,dur=600)
        bar=rect(slide,Inches(0),Inches(0),Inches(0.6),H,rgb(th["accent"]))
        seq.wipe_in(bar,dur=400)
        t=txt(slide,data["title"],Inches(0.8),Inches(0.18),Inches(11),Inches(1.0),
              th["font_head"],36,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=200,from_dir="left")
        bot=rect(slide,Inches(0.8),Inches(1.22),Inches(10),Inches(0.07),rgb(th["accent"]))
        seq.wipe_in(bot,dur=500,delay=350)
        pts=data.get("points",[""])
        pos=[(Inches(0.8),Inches(1.48)),(Inches(6.9),Inches(1.48)),(Inches(0.8),Inches(4.18)),(Inches(6.9),Inches(4.18))]
        for i,(pt,(cx,cy)) in enumerate(zip(pts[:4],pos)):
            d=440+i*140
            c=rect(slide,cx,cy,Inches(5.85),Inches(2.5),rgb(th["card_bg"]))
            seq.fade(c,dur=400,delay=d)
            at=rect(slide,cx,cy,Inches(0.5),Inches(2.5),rgb(th["accent"]))
            seq.wipe_in(at,dur=350,delay=d+60)
            num=txt(slide,f"{i+1}",cx,cy+Inches(0.88),Inches(0.5),Inches(0.7),
                    th["font_head"],20,rgb(th["bg"]),bold=True,align=PP_ALIGN.CENTER)
            seq.appear(num,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.64),cy+Inches(0.08),Inches(5.1),Inches(2.28),
                     th["font_body"],13,rgb(th["body_color"]),wrap=True)
            seq.fly_in(ptxt,dur=400,delay=d+150,from_dir="bottom")
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx): return self.grid(prs,data,idx)
    def timeline(self,prs,data,idx): return self.content(prs,data,idx)


# ════════════════════════════════════════════════════════
# T8: NEON NIGHTS — outlined shapes, synthwave, glowing
# ════════════════════════════════════════════════════════
class NeonNights(_Base):
    trans="flip"
    def _outline(self,slide,l,t,w,h,col,seq,d=0):
        s=slide.shapes.add_shape(1,l,t,w,h)
        s.fill.background(); s.line.color.rgb=col; s.line.width=Pt(1.5)
        seq.zoom_in(s,dur=350,delay=d); return s

    def hero(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        for r_i,sz in enumerate([Inches(8),Inches(6),Inches(4)]):
            cx=(W-sz)/2; cy=(H-sz)/2
            c=slide.shapes.add_shape(9,cx,cy,sz,sz)
            c.fill.background(); c.line.color.rgb=rgb(th["accent"]); c.line.width=Pt(0.5+r_i*0.3)
            seq.zoom_in(c,dur=500+r_i*100,delay=r_i*100)
            if r_i==1: seq.pulse(c)
        bar=self._outline(slide,Inches(0.3),Inches(0.18),W-Inches(0.6),Inches(0.04),rgb(th["accent"]),seq,300)
        n=txt(slide,f"{idx:02d}",Inches(0.5),Inches(0.28),Inches(1.2),Inches(0.5),
              th["font_body"],11,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=350)
        t=txt(slide,data["title"],Inches(0.5),Inches(1.38),Inches(12.3),Inches(3.0),
              th["font_head"],56,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=600,delay=350)
        s=txt(slide,data.get("subtitle",""),Inches(0.5),Inches(4.28),Inches(11),Inches(0.7),
              th["font_body"],18,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=400,delay=600)
        d=txt(slide,data.get("detail",""),Inches(0.5),Inches(5.18),Inches(11),Inches(2.0),
              th["font_body"],14,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=400,delay=800)
        return self._fin(slide,seq,data)

    def content(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        bar=rect(slide,Inches(0),Inches(0),Inches(0.4),H,rgb(th["accent"]))
        seq.wipe_in(bar,dur=400)
        hbox=self._outline(slide,Inches(0.5),Inches(0.08),Inches(12.5),Inches(1.42),rgb(th["accent"]),seq,100)
        n=txt(slide,f"{idx:02d}",Inches(0.65),Inches(0.18),Inches(1.2),Inches(0.5),
              th["font_body"],11,rgb(th["accent"]),bold=True)
        seq.appear(n,delay=200)
        t=txt(slide,data["title"],Inches(0.65),Inches(0.38),Inches(12),Inches(1.0),
              th["font_head"],34,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=250)
        s=txt(slide,data.get("subtitle",""),Inches(0.65),Inches(1.53),Inches(7.5),Inches(0.5),
              th["font_body"],13,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=350,delay=420)
        bl=bullets(slide,data.get("points",[]),Inches(0.65),Inches(2.18),Inches(6.5),Inches(4.8),th,bullet="⟩  ")
        seq.fly_in(bl,dur=500,delay=500)
        rc=self._outline(slide,Inches(7.35),Inches(0.08),Inches(5.65),Inches(7.05),rgb(th["accent2"]),seq,300)
        seq.pulse(rc)
        d=txt(slide,data.get("detail",""),Inches(7.55),Inches(0.68),Inches(5.25),Inches(6.2),
              th["font_body"],14,rgb(th["body_color"]),wrap=True)
        seq.fade(d,dur=400,delay=550)
        return self._fin(slide,seq,data)

    def grid(self,prs,data,idx):
        slide,seq=self._slide(prs); th=self.th
        bar=rect(slide,Inches(0),Inches(0),Inches(0.4),H,rgb(th["accent"]))
        seq.wipe_in(bar,dur=400)
        t=txt(slide,data["title"],Inches(0.65),Inches(0.18),Inches(12),Inches(1.0),
              th["font_head"],36,rgb(th["title_color"]),bold=True)
        seq.fly_in(t,dur=500,delay=200)
        s=txt(slide,data.get("subtitle",""),Inches(0.65),Inches(1.18),Inches(10),Inches(0.45),
              th["font_body"],13,rgb(th["accent"]),italic=True)
        seq.fade(s,dur=350,delay=350)
        pts=data.get("points",[""])
        pos=[(Inches(0.65),Inches(1.82)),(Inches(7.0),Inches(1.82)),(Inches(0.65),Inches(4.52)),(Inches(7.0),Inches(4.52))]
        for i,(pt,(cx,cy)) in enumerate(zip(pts[:4],pos)):
            d=440+i*130
            c=self._outline(slide,cx,cy,Inches(6.1),Inches(2.45),rgb(th["accent"]),seq,d)
            if i%2==1: seq.pulse(c)
            nc=rect(slide,cx+Inches(0.14),cy+Inches(0.12),Inches(0.55),Inches(0.55),rgb(th["accent"]))
            seq.zoom_in(nc,dur=250,delay=d+80)
            bn=txt(slide,f"0{i+1}",cx+Inches(0.14),cy+Inches(0.12),Inches(0.55),Inches(0.55),
                   th["font_body"],12,rgb(th["bg"]),bold=True,align=PP_ALIGN.CENTER)
            seq.appear(bn,delay=d+80)
            ptxt=txt(slide,pt,cx+Inches(0.84),cy+Inches(0.08),Inches(5.1),Inches(2.25),
                     th["font_body"],12,rgb(th["body_color"]),wrap=True)
            seq.fly_in(ptxt,dur=380,delay=d+150,from_dir="bottom")
        return self._fin(slide,seq,data)

    def numbered(self,prs,data,idx): return self.grid(prs,data,idx)
    def timeline(self,prs,data,idx): return self.content(prs,data,idx)


# ════════════════════════════════════════════════════════
# Registry + builder
# ════════════════════════════════════════════════════════

TEMPLATES = {
    "futuristic":  Futuristic,
    "minimalist":  Minimalist,
    "artistic":    Artistic,
    "misty":       MistyPeaks,
    "monochrome":  Monochrome,
    "hardware":    Hardware,
    "magazine":    Magazine,
    "neon":        NeonNights,
}

LAYOUT_CYCLE = ["content","grid","numbered","timeline","content","numbered","grid","timeline","content","grid"]
LAYOUT_MAP   = {
    "title_hero":"hero","two_column":"content","icon_grid":"grid",
    "stat_callout":"numbered","timeline":"timeline","full_detail":"numbered",
}

def build_ppt(outline, theme, style="futuristic"):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    Cls = TEMPLATES.get(style, Futuristic)
    tmpl = Cls(theme)
    n = len(outline)
    for idx, data in enumerate(outline, start=1):
        if idx==1 or idx==n:
            tmpl.hero(prs, data, idx)
        else:
            key = LAYOUT_MAP.get(data.get("layout",""))
            if not key:
                key = LAYOUT_CYCLE[(idx-2) % len(LAYOUT_CYCLE)]
            getattr(tmpl, key, tmpl.content)(prs, data, idx)
    buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf.read()